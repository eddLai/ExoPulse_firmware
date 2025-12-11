#!/usr/bin/env python3
"""
Script to convert uMyo_Blueprint_to_Beacon PDF to PowerPoint presentation
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

def create_presentation():
    # Create presentation with 16:9 aspect ratio
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    # Define standard layout
    blank_slide_layout = prs.slide_layouts[6]  # Blank layout

    # Slide 1: Title Slide
    slide = prs.slides.add_slide(blank_slide_layout)

    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = "uMyo: From Blueprint to Beacon"
    title_para.font.size = Pt(48)
    title_para.font.bold = True
    title_para.alignment = PP_ALIGN.CENTER

    # Add subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(9), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.text = "An Engineering Journey Through the nRF SoC-Based EMG/IMU Sensor"
    subtitle_para.font.size = Pt(24)
    subtitle_para.alignment = PP_ALIGN.CENTER

    # Slide 2: The Heart of the System
    slide = prs.slides.add_slide(blank_slide_layout)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = "The Heart of the System: uMyo Wireless EMG/IMU Sensor"
    title_para.font.size = Pt(32)
    title_para.font.bold = True

    # Body text
    body_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(4.5), Inches(4))
    body_frame = body_box.text_frame
    body_frame.word_wrap = True

    # Body Text section
    p = body_frame.paragraphs[0]
    p.text = "Body Text:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.space_after = Pt(8)

    # Bullet 1
    p = body_frame.add_paragraph()
    p.text = "What is uMyo?: A wearable sensor for capturing Electromyography (EMG) and Inertial Measurement Unit (IMU) data."
    p.level = 0
    p.font.size = Pt(12)
    p.space_after = Pt(6)

    # Bullet 2
    p = body_frame.add_paragraph()
    p.text = "Primary Role: Serves as the recommended wireless signal acquisition method for the ExoPulse exoskeleton control system."
    p.level = 0
    p.font.size = Pt(12)
    p.space_after = Pt(6)

    # Bullet 3
    p = body_frame.add_paragraph()
    p.text = "Core Function: Operates as a BLE Beacon, continuously broadcasting sensor data (EMG, IMU, battery status) in a one-way, connectionless stream."
    p.level = 0
    p.font.size = Pt(12)
    p.space_after = Pt(12)

    # Hardware Specifications
    p = body_frame.add_paragraph()
    p.text = "Key Hardware Specifications:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.space_after = Pt(8)

    p = body_frame.add_paragraph()
    p.text = "MCU: nRF52810 (BLE System-on-Chip)"
    p.level = 0
    p.font.size = Pt(12)
    p.space_after = Pt(4)

    p = body_frame.add_paragraph()
    p.text = "Sensors: LSM6DS3 (6-DOF Accel + Gyro), QMC7983 (Magnetometer)"
    p.level = 0
    p.font.size = Pt(12)
    p.space_after = Pt(4)

    p = body_frame.add_paragraph()
    p.text = "Power: ~180mAh Li-Po battery providing an estimated 10-15 hours of continuous operation."
    p.level = 0
    p.font.size = Pt(12)
    p.space_after = Pt(4)

    p = body_frame.add_paragraph()
    p.text = "EMG Sampling: Single channel ADC at a 1kHz sample rate."
    p.level = 0
    p.font.size = Pt(12)

    # Slide 3: Hardware Architecture
    slide = prs.slides.add_slide(blank_slide_layout)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = "uMyo v3.1 Hardware Architecture"
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.alignment = PP_ALIGN.CENTER

    # Architecture description
    desc_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(3.5))
    desc_frame = desc_box.text_frame
    desc_frame.word_wrap = True

    p = desc_frame.paragraphs[0]
    p.text = "System Architecture:"
    p.font.size = Pt(20)
    p.font.bold = True
    p.space_after = Pt(10)

    p = desc_frame.add_paragraph()
    p.text = "• EMG Signal Chain: High-Pass Filter → Low-Pass (Anti-Aliasing) Filter → Gain Stage → nRF52810 ADC (P0.28/AIN4) for 1kHz sampling"
    p.font.size = Pt(14)
    p.space_after = Pt(8)

    p = desc_frame.add_paragraph()
    p.text = "• CPU & Radio Core: nRF52810-QFAA MCU with 32MHz main crystal and 32.768kHz RTC crystal"
    p.font.size = Pt(14)
    p.space_after = Pt(8)

    p = desc_frame.add_paragraph()
    p.text = "• IMU & Sensors: LSM6DS3 (Accel + Gyro) and QMC5883L Magnetometer connected via I2C Bus"
    p.font.size = Pt(14)
    p.space_after = Pt(8)

    p = desc_frame.add_paragraph()
    p.text = "• RF Output: PCB Trace Antenna (AE1) with matching network for 50Ω impedance"
    p.font.size = Pt(14)
    p.space_after = Pt(8)

    p = desc_frame.add_paragraph()
    p.text = "• Power Management: 180mAh Li-Po Battery → XC6206P332MR (LDO) → Regulated 3.3V"
    p.font.size = Pt(14)
    p.space_after = Pt(8)

    p = desc_frame.add_paragraph()
    p.text = "• User Interface: RGB LED, Mode Button, SWD Interface (SWDIO/SWCLK)"
    p.font.size = Pt(14)

    # Slide 4: Radio Protocols
    slide = prs.slides.add_slide(blank_slide_layout)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = "A Sensor with Three Voices: The uMyo Radio Protocols"
    title_para.font.size = Pt(30)
    title_para.font.bold = True
    title_para.alignment = PP_ALIGN.CENTER

    # Intro text
    intro_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(9), Inches(0.5))
    intro_frame = intro_box.text_frame
    intro_para = intro_frame.paragraphs[0]
    intro_para.text = "uMyo supports three transmission modes, selectable via a button press and indicated by the RGB LED. This presentation focuses on the universally compatible BLE mode."
    intro_para.font.size = Pt(13)
    intro_para.alignment = PP_ALIGN.CENTER
    intro_frame.word_wrap = True

    # Table
    table_data = [
        ["Mode", "LED Indicator", "Protocol", "Data Rate", "Best For", "Receiver Required"],
        ["BLE", "Blue\n(3 pulses)", "BLE\nAdvertisement", "~10-100 Hz", "Universal\ncompatibility", "ESP32,\nSmartphones, PC"],
        ["Fast32", "Magenta\n(3 pulses)", "Nordic\nProprietary", "~1000 Hz", "High-speed\nsingle sensor", "nRF52 device"],
        ["Fast64", "Green\n(3 pulses)", "Custom TDMA\nStar", "~500 Hz per\nsensor", "Synchronized\nmulti-sensor", "nRF52 central\nhub"]
    ]

    # Add table
    rows = len(table_data)
    cols = len(table_data[0])
    left = Inches(0.5)
    top = Inches(1.8)
    width = Inches(9)
    height = Inches(3)

    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # Set column widths
    table.columns[0].width = Inches(1.2)
    table.columns[1].width = Inches(1.3)
    table.columns[2].width = Inches(1.5)
    table.columns[3].width = Inches(1.5)
    table.columns[4].width = Inches(1.8)
    table.columns[5].width = Inches(1.7)

    # Fill table
    for i, row_data in enumerate(table_data):
        for j, cell_text in enumerate(row_data):
            cell = table.cell(i, j)
            cell.text = cell_text

            # Format text
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(11)
                paragraph.alignment = PP_ALIGN.CENTER

                if i == 0:  # Header row
                    paragraph.font.bold = True
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(220, 220, 220)

            cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Slide 5: JLCPCB Production Log
    slide = prs.slides.add_slide(blank_slide_layout)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = "From Design to Reality: A JLCPCB Production Log"
    title_para.font.size = Pt(34)
    title_para.font.bold = True
    title_para.alignment = PP_ALIGN.CENTER

    # Body text
    body_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(3.5))
    body_frame = body_box.text_frame
    body_frame.word_wrap = True

    p = body_frame.paragraphs[0]
    p.text = "The uMyo v3.1 was designed in KiCad 6 and sent for SMT assembly (PCBA) at JLCPCB. The process revealed several critical lessons in design-for-manufacturing (DFM) and supply chain management. The following challenges document this journey."
    p.font.size = Pt(18)
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(20)

    p = body_frame.add_paragraph()
    p.text = "Key Lessons Learned:"
    p.font.size = Pt(20)
    p.font.bold = True
    p.space_after = Pt(10)

    p = body_frame.add_paragraph()
    p.text = "1. Missing Drill Layer - File verification is critical"
    p.font.size = Pt(16)
    p.level = 0
    p.space_after = Pt(6)

    p = body_frame.add_paragraph()
    p.text = "2. Component Shortages - De-scope non-critical features"
    p.font.size = Pt(16)
    p.level = 0
    p.space_after = Pt(6)

    p = body_frame.add_paragraph()
    p.text = "3. DFM and RF Verification - Engineering review prevented critical failures"
    p.font.size = Pt(16)
    p.level = 0

    # Slide 6: Production Challenge 1
    slide = prs.slides.add_slide(blank_slide_layout)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = "Production Challenge 1: The Missing Drill Layer"
    title_para.font.size = Pt(32)
    title_para.font.bold = True

    # Content
    body_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(9), Inches(4))
    body_frame = body_box.text_frame
    body_frame.word_wrap = True

    p = body_frame.paragraphs[0]
    run = p.add_run()
    run.text = "✱ The Problem: "
    run.font.size = Pt(16)
    run.font.bold = True
    p.space_after = Pt(4)

    run = p.add_run()
    run.text = 'The initial order was rejected after upload with the message "Missing Drill Layer." The .drl file was accidentally omitted from the Gerber ZIP archive.'
    run.font.size = Pt(14)

    p = body_frame.add_paragraph()
    run = p.add_run()
    run.text = "✱ The Complication: "
    run.font.size = Pt(16)
    run.font.bold = True
    p.space_after = Pt(4)

    run = p.add_run()
    run.text = "JLCPCB's system does not allow direct file replacement on an existing order for integrity reasons. Contacting support could not resolve the issue."
    run.font.size = Pt(14)
    p.space_after = Pt(10)

    p = body_frame.add_paragraph()
    run = p.add_run()
    run.text = "✱ The Solution: "
    run.font.size = Pt(16)
    run.font.bold = True
    p.space_after = Pt(4)

    run = p.add_run()
    run.text = "The original order had to be cancelled and refunded. A new order was placed with a correctly packaged ZIP file containing all layers, including the drill file."
    run.font.size = Pt(14)
    p.space_after = Pt(14)

    # Lesson box
    p = body_frame.add_paragraph()
    run = p.add_run()
    run.text = "✱ The Lesson: "
    run.font.size = Pt(16)
    run.font.bold = True
    p.space_after = Pt(4)

    run = p.add_run()
    run.text = "Always verify the contents of your manufacturing ZIP archive before upload. KiCad's workflow can make it easy to miss this step, leading to a 1-2 day delay in the review process."
    run.font.size = Pt(14)

    # Add colored box around lesson
    lesson_shape = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0.5), Inches(3.8), Inches(9), Inches(1.2)
    )
    lesson_shape.fill.solid()
    lesson_shape.fill.fore_color.rgb = RGBColor(220, 240, 255)
    lesson_shape.line.color.rgb = RGBColor(100, 150, 200)
    lesson_shape.line.width = Pt(2)

    # Move text box to front
    slide.shapes._spTree.remove(body_box._element)
    slide.shapes._spTree.append(body_box._element)

    # Slide 7: Production Challenge 2
    slide = prs.slides.add_slide(blank_slide_layout)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = "Production Challenge 2: Navigating Component Shortages"
    title_para.font.size = Pt(30)
    title_para.font.bold = True

    # Content
    body_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(5.5), Inches(4))
    body_frame = body_box.text_frame
    body_frame.word_wrap = True

    p = body_frame.paragraphs[0]
    p.text = "• The Problem: "
    p.font.size = Pt(15)
    p.font.bold = True
    p.space_after = Pt(4)

    run = p.add_run()
    run.text = "During BOM review, two key components were found to have zero stock at JLCPCB:"
    run.font.size = Pt(13)

    p = body_frame.add_paragraph()
    p.text = "○ U6: QMC5883L Magnetometer"
    p.font.size = Pt(12)
    p.level = 1
    p.space_after = Pt(4)

    p = body_frame.add_paragraph()
    p.text = "○ D1: RGB LED"
    p.font.size = Pt(12)
    p.level = 1
    p.space_after = Pt(10)

    p = body_frame.add_paragraph()
    p.text = "• The Vendor's Option: "
    p.font.size = Pt(15)
    p.font.bold = True
    p.space_after = Pt(4)

    run = p.add_run()
    run.text = 'Use "Global Sourcing" to procure the parts from the worldwide market, a process known to cause significant delays.'
    run.font.size = Pt(13)
    p.space_after = Pt(10)

    p = body_frame.add_paragraph()
    p.text = "• The Strategic Decision: "
    p.font.size = Pt(15)
    p.font.bold = True
    p.space_after = Pt(4)

    run = p.add_run()
    run.text = 'For a first-revision prototype, verifying core functionality is the top priority. To avoid stalling the entire project, both components were marked "Do Not Populate" (DNP).'
    run.font.size = Pt(13)
    p.space_after = Pt(12)

    # Lesson box
    lesson_shape = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0.5), Inches(3.9), Inches(5.5), Inches(1.3)
    )
    lesson_shape.fill.solid()
    lesson_shape.fill.fore_color.rgb = RGBColor(255, 250, 220)
    lesson_shape.line.color.rgb = RGBColor(200, 150, 100)
    lesson_shape.line.width = Pt(2)

    lesson_box = slide.shapes.add_textbox(Inches(0.6), Inches(4), Inches(5.3), Inches(1.1))
    lesson_frame = lesson_box.text_frame
    lesson_frame.word_wrap = True

    p = lesson_frame.paragraphs[0]
    p.text = "The Lesson: "
    p.font.size = Pt(14)
    p.font.bold = True
    p.space_after = Pt(4)

    run = p.add_run()
    run.text = "In prototyping, be decisive about de-scoping non-critical features. Don't let a component shortage derail your validation schedule. Get the board in hand, test the core, and add the rest later."
    run.font.size = Pt(12)

    # Slide 8: Production Challenge 3
    slide = prs.slides.add_slide(blank_slide_layout)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = "Production Challenge 3: DFM and RF Polarity Verification"
    title_para.font.size = Pt(28)
    title_para.font.bold = True

    # Content
    body_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(9), Inches(4))
    body_frame = body_box.text_frame
    body_frame.word_wrap = True

    p = body_frame.paragraphs[0]
    p.text = "• The Process: "
    p.font.size = Pt(15)
    p.font.bold = True
    p.space_after = Pt(4)

    run = p.add_run()
    run.text = "JLCPCB's engineering review team (contact: Doris) provided a thorough audit of the design before production."
    run.font.size = Pt(13)
    p.space_after = Pt(8)

    p = body_frame.add_paragraph()
    p.text = "• The Critical Question: "
    p.font.size = Pt(15)
    p.font.bold = True
    p.space_after = Pt(4)

    run = p.add_run()
    run.text = "The engineers flagged the chip antenna (AE1) and requested verification of its placement orientation and polarity. This is crucial for RF performance."
    run.font.size = Pt(13)
    p.space_after = Pt(8)

    p = body_frame.add_paragraph()
    p.text = "• The Verification:"
    p.font.size = Pt(15)
    p.font.bold = True
    p.space_after = Pt(6)

    p = body_frame.add_paragraph()
    p.text = "1. Cross-referenced the KiCad PCB layout with the antenna's component datasheet."
    p.font.size = Pt(13)
    p.level = 0
    p.space_after = Pt(4)

    p = body_frame.add_paragraph()
    p.text = "2. Confirmed that Pad 1 is the Feed pin."
    p.font.size = Pt(13)
    p.level = 0
    p.space_after = Pt(4)

    p = body_frame.add_paragraph()
    p.text = "3. Verified that the current layout orientation was correct."
    p.font.size = Pt(13)
    p.level = 0
    p.space_after = Pt(8)

    p = body_frame.add_paragraph()
    p.text = "• The Outcome: "
    p.font.size = Pt(15)
    p.font.bold = True
    p.space_after = Pt(4)

    run = p.add_run()
    run.text = "After providing confirmation, the engineering hold was released, and production proceeded. This step prevented a potentially critical performance flaw in the final boards."
    run.font.size = Pt(13)

    # Slide 9: BLE Advertisement
    slide = prs.slides.add_slide(blank_slide_layout)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = "The Language of Beacons: Why uMyo Uses BLE Advertisement"
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.alignment = PP_ALIGN.CENTER

    # Intro
    intro_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(9), Inches(0.5))
    intro_frame = intro_box.text_frame
    intro_frame.word_wrap = True
    p = intro_frame.paragraphs[0]
    p.text = "BLE Advertisement is a one-way broadcast mechanism. A device sends small data packets to any nearby receiver without establishing a formal connection, much like a radio station."
    p.font.size = Pt(14)
    p.alignment = PP_ALIGN.CENTER

    # Table
    table_data = [
        ["Feature", "Advertisement Mode (uMyo)", "Connection Mode"],
        ["Power Consumption", "Very Low", "Higher"],
        ["Multiple Receivers", "Yes (unlimited)", "No (one at a time)"],
        ["Setup Required", "None (Power-on & go)", "Pairing / Bonding"],
        ["Use Case", "Broadcasting sensor data", "Streaming data, control"]
    ]

    # Add table
    rows = len(table_data)
    cols = len(table_data[0])
    left = Inches(0.8)
    top = Inches(1.9)
    width = Inches(8.4)
    height = Inches(2)

    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # Set column widths
    table.columns[0].width = Inches(2.5)
    table.columns[1].width = Inches(3.2)
    table.columns[2].width = Inches(2.7)

    # Fill table
    for i, row_data in enumerate(table_data):
        for j, cell_text in enumerate(row_data):
            cell = table.cell(i, j)
            cell.text = cell_text

            # Format text
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(13)
                paragraph.alignment = PP_ALIGN.CENTER

                if i == 0:  # Header row
                    paragraph.font.bold = True
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(220, 220, 220)

            cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Key Advantages
    adv_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1.2))
    adv_frame = adv_box.text_frame
    adv_frame.word_wrap = True

    p = adv_frame.paragraphs[0]
    p.text = "Key Advantages for uMyo:"
    p.font.size = Pt(16)
    p.font.bold = True
    p.space_after = Pt(6)

    p = adv_frame.add_paragraph()
    p.text = "✱ Multiple ESP32s can receive data simultaneously."
    p.font.size = Pt(13)
    p.space_after = Pt(3)

    p = adv_frame.add_paragraph()
    p.text = "✱ Extremely simple operation with no pairing needed."
    p.font.size = Pt(13)
    p.space_after = Pt(3)

    p = adv_frame.add_paragraph()
    p.text = "✱ Optimized for long battery life in a wearable device."
    p.font.size = Pt(13)

    # Slide 10: Anatomy of Broadcast Packet
    slide = prs.slides.add_slide(blank_slide_layout)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = "Anatomy of a uMyo Broadcast Packet"
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.alignment = PP_ALIGN.CENTER

    # Intro
    intro_box = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(9), Inches(0.4))
    intro_frame = intro_box.text_frame
    intro_frame.word_wrap = True
    p = intro_frame.paragraphs[0]
    p.text = "Each BLE advertisement packet is limited to a 31-byte payload. uMyo uses a specific packet type and structure to broadcast its data efficiently."
    p.font.size = Pt(13)
    p.alignment = PP_ALIGN.CENTER

    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(4.5), Inches(3.8))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True

    p = left_frame.paragraphs[0]
    p.text = "Packet Type: ADV_NONCONN_IND"
    p.font.size = Pt(14)
    p.font.bold = True
    p.space_after = Pt(6)

    p = left_frame.add_paragraph()
    run = p.add_run()
    run.text = "✱ This signifies a "
    run.font.size = Pt(12)
    run = p.add_run()
    run.text = "Non-Connectable, Non-Scannable"
    run.font.size = Pt(12)
    run.font.bold = True
    run = p.add_run()
    run.text = ' broadcast. uMyo is a "shout-only" device; it will not respond to connection or scan requests, maximizing efficiency.'
    run.font.size = Pt(12)
    p.space_after = Pt(12)

    p = left_frame.add_paragraph()
    p.text = "AdvData (Payload) Structure:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.space_after = Pt(6)

    p = left_frame.add_paragraph()
    p.text = "✱ The payload is composed of one or more 'AD Structures.'"
    p.font.size = Pt(12)
    p.space_after = Pt(6)

    p = left_frame.add_paragraph()
    p.text = "✱ Each structure has a [Length] [Type] [Data] format."
    p.font.size = Pt(12)
    p.space_after = Pt(12)

    p = left_frame.add_paragraph()
    p.text = "uMyo's Key AD Structures:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.space_after = Pt(6)

    p = left_frame.add_paragraph()
    run = p.add_run()
    run.text = "✱ 0x08 - Short Name: "
    run.font.size = Pt(12)
    run.font.bold = True
    run = p.add_run()
    run.text = 'Contains the device name, e.g., "uMyo v2".'
    run.font.size = Pt(12)
    p.space_after = Pt(6)

    p = left_frame.add_paragraph()
    run = p.add_run()
    run.text = "✱ 0xFF - Manufacturer Data: "
    run.font.size = Pt(12)
    run.font.bold = True
    run = p.add_run()
    run.text = "The primary container for custom sensor data (EMG, IMU, battery level), formatted according to a proprietary structure."
    run.font.size = Pt(12)

    # Slide 11: ESP32 vs PC Scanning
    slide = prs.slides.add_slide(blank_slide_layout)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = "An Expert's Insight: Why Can My ESP32 See uMyo, but My PC Can't?"
    title_para.font.size = Pt(24)
    title_para.font.bold = True
    title_para.alignment = PP_ALIGN.CENTER

    # Table
    table_data = [
        ["Scanner", "uMyo Detected?", "Observation"],
        ["ESP32 BLE Scanner", "✓ Yes", "Reliably finds device\nd7:72:58:3a:82:d4"],
        ["PC (Linux bluetoothctl)", "✗ No", "Device does not appear\nin scan results"]
    ]

    # Add table
    rows = len(table_data)
    cols = len(table_data[0])
    left = Inches(1.5)
    top = Inches(1.3)
    width = Inches(7)
    height = Inches(1.5)

    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # Set column widths
    table.columns[0].width = Inches(2.5)
    table.columns[1].width = Inches(2)
    table.columns[2].width = Inches(2.5)

    # Fill table
    for i, row_data in enumerate(table_data):
        for j, cell_text in enumerate(row_data):
            cell = table.cell(i, j)
            cell.text = cell_text

            # Format text
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(14)
                paragraph.alignment = PP_ALIGN.CENTER

                if i == 0:  # Header row
                    paragraph.font.bold = True
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(220, 220, 220)

            cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Explanation
    exp_box = slide.shapes.add_textbox(Inches(0.8), Inches(3.2), Inches(8.4), Inches(2))
    exp_frame = exp_box.text_frame
    exp_frame.word_wrap = True

    p = exp_frame.paragraphs[0]
    p.text = "This behavior is not a fault. It is a direct result of uMyo's specific BLE configuration interacting with the default filtering mechanisms of most PC Bluetooth stacks. The following slides explain the root causes."
    p.font.size = Pt(16)
    p.alignment = PP_ALIGN.CENTER

    # Slide 12: Root Cause 1
    slide = prs.slides.add_slide(blank_slide_layout)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = "Root Cause 1: Filtering of Non-Connectable Devices"
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.alignment = PP_ALIGN.CENTER

    # Content
    body_box = slide.shapes.add_textbox(Inches(5.5), Inches(1.2), Inches(4), Inches(4))
    body_frame = body_box.text_frame
    body_frame.word_wrap = True

    p = body_frame.paragraphs[0]
    p.text = "• uMyo's Behavior: "
    p.font.size = Pt(14)
    p.font.bold = True
    p.space_after = Pt(4)

    run = p.add_run()
    run.text = 'uMyo broadcasts using the "ADV_NONCONN_IND" type. It explicitly states, "I cannot be connected to."'
    run.font.size = Pt(12)
    p.space_after = Pt(10)

    p = body_frame.add_paragraph()
    p.text = "• PC Stack's Behavior: "
    p.font.size = Pt(14)
    p.font.bold = True
    p.space_after = Pt(4)

    run = p.add_run()
    run.text = "Most PC Bluetooth software (like Linux's BlueZ) is designed to find peripherals you can connect to (keyboards, headphones, mice). By default, they often "
    run.font.size = Pt(12)
    run = p.add_run()
    run.text = "filter out and hide non-connectable devices"
    run.font.size = Pt(12)
    run.font.bold = True
    run = p.add_run()
    run.text = " to de-clutter the user interface."
    run.font.size = Pt(12)
    p.space_after = Pt(10)

    p = body_frame.add_paragraph()
    p.text = "• The Result: "
    p.font.size = Pt(14)
    p.font.bold = True
    p.space_after = Pt(4)

    run = p.add_run()
    run.text = "The PC sees the uMyo packet but chooses to ignore it as irrelevant for its primary purpose."
    run.font.size = Pt(12)

    # Slide 13: Root Cause 2
    slide = prs.slides.add_slide(blank_slide_layout)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = "Root Cause 2: The Absence of a Scan Response"
    title_para.font.size = Pt(34)
    title_para.font.bold = True
    title_para.alignment = PP_ALIGN.CENTER

    # Content
    body_box = slide.shapes.add_textbox(Inches(5.5), Inches(1.2), Inches(4), Inches(4))
    body_frame = body_box.text_frame
    body_frame.word_wrap = True

    p = body_frame.paragraphs[0]
    p.text = "• Active vs. Passive Scanning:"
    p.font.size = Pt(14)
    p.font.bold = True
    p.space_after = Pt(6)

    p = body_frame.add_paragraph()
    p.text = "○ An 'active' scanner sends a SCAN_REQ packet to ask for more information (like a full device name)."
    p.font.size = Pt(11)
    p.level = 1
    p.space_after = Pt(4)

    p = body_frame.add_paragraph()
    p.text = "○ A 'passive' scanner just listens to the initial advertisement broadcasts."
    p.font.size = Pt(11)
    p.level = 1
    p.space_after = Pt(10)

    p = body_frame.add_paragraph()
    p.text = "• uMyo's Behavior: "
    p.font.size = Pt(14)
    p.font.bold = True
    p.space_after = Pt(4)

    run = p.add_run()
    run.text = "Because it is non-scannable (ADV_NONCONN_IND), uMyo is designed to "
    run.font.size = Pt(11)
    run = p.add_run()
    run.text = "completely ignore"
    run.font.size = Pt(11)
    run.font.bold = True
    run = p.add_run()
    run.text = " `SCAN_REQ` packets. It never sends a `SCAN_RSP` (Scan Response)."
    run.font.size = Pt(11)
    p.space_after = Pt(10)

    p = body_frame.add_paragraph()
    p.text = "• PC Stack's Behavior: "
    p.font.size = Pt(14)
    p.font.bold = True
    p.space_after = Pt(4)

    run = p.add_run()
    run.text = "Many PC tools perform an active scan and rely on the `SCAN_RSP` to populate their device list. When no response is received, they may fail to list the device correctly."
    run.font.size = Pt(11)

    # Slide 14: ESP32 Advantage
    slide = prs.slides.add_slide(blank_slide_layout)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = "The ESP32's Advantage: Raw, Unfiltered Scanning"
    title_para.font.size = Pt(34)
    title_para.font.bold = True
    title_para.alignment = PP_ALIGN.CENTER

    # Content - two columns
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.5), Inches(4))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True

    p = left_frame.paragraphs[0]
    p.text = "• ESP32's Method: "
    p.font.size = Pt(14)
    p.font.bold = True
    p.space_after = Pt(6)

    run = p.add_run()
    run.text = "The ESP32's BLE libraries allow for raw, low-level scanning. They capture and report *every* advertisement packet seen on the air, without applying the high-level filters common on PCs."
    run.font.size = Pt(12)

    right_box = slide.shapes.add_textbox(Inches(5.5), Inches(1.2), Inches(4), Inches(4))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True

    p = right_frame.paragraphs[0]
    p.text = "• Payload-First Parsing: "
    p.font.size = Pt(14)
    p.font.bold = True
    p.space_after = Pt(6)

    run = p.add_run()
    run.text = "The ESP32 code doesn't rely on specific AD types or scan responses. The `uMyo_BLE` library scans the raw payload bytes directly for the `uMyo` name and manufacturer data."
    run.font.size = Pt(12)
    p.space_after = Pt(10)

    p = right_frame.add_paragraph()
    p.text = "• Conclusion: "
    p.font.size = Pt(14)
    p.font.bold = True
    p.space_after = Pt(6)

    run = p.add_run()
    run.text = "The ESP32 is the superior tool for this task because it acts as a true protocol analyzer, not a consumer peripheral manager. It sees what is actually there, not what a higher-level stack thinks is relevant."
    run.font.size = Pt(12)

    # Slide 15: Recommended Receivers
    slide = prs.slides.add_slide(blank_slide_layout)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = "Recommended Receivers and Fast64 Data Reference"
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.alignment = PP_ALIGN.CENTER

    # Left column - Recommendations
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(4.5), Inches(4.2))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True

    p = left_frame.paragraphs[0]
    p.text = "For reliable data acquisition, choose the receiver that matches your required protocol and data rate."
    p.font.size = Pt(13)
    p.space_after = Pt(12)

    p = left_frame.add_paragraph()
    run = p.add_run()
    run.text = "• For BLE Mode: An "
    run.font.size = Pt(12)
    run = p.add_run()
    run.text = "ESP32-based scanner"
    run.font.size = Pt(12)
    run.font.bold = True
    run = p.add_run()
    run.text = " is recommended for its raw scanning capabilities. Standard PC tools may filter the uMyo's non-connectable advertisements."
    run.font.size = Pt(12)
    p.space_after = Pt(10)

    p = left_frame.add_paragraph()
    run = p.add_run()
    run.text = "• For Fast32/Fast64 Mode: An "
    run.font.size = Pt(12)
    run = p.add_run()
    run.text = "nRF52840 DK"
    run.font.size = Pt(12)
    run.font.bold = True
    run = p.add_run()
    run.text = " (or similar nRF52 device) is required to run the proprietary high-speed protocols."
    run.font.size = Pt(12)

    # Right column - Fast64 Table
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.1), Inches(4.3), Inches(0.6))
    right_frame = right_box.text_frame
    p = right_frame.paragraphs[0]
    p.text = "Reference: Fast64 Protocol Packet Structure\n(62-byte Payload)"
    p.font.size = Pt(13)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Table
    table_data = [
        ["Offset", "Size", "Field", "Description"],
        ["0", "1", "Packet_ID", "Incrementing counter"],
        ["2-5", "4", "Unit_ID", "Unique sensor identifier"],
        ["8", "1", "Battery", "Battery level (0-255)"],
        ["11-26", "16", "ADC Samples [8]", "8 consecutive 16-bit\nsamples @ 1kHz"],
        ["35-42", "8", "Quaternion [4]", "w, x, y, z (16-bit each)"],
        ["43-48", "6", "Accelerometer [3]", "X, Y, Z acceleration"],
        ["49-54", "6", "YPR [3]", "Yaw, Pitch, Roll"],
        ["55-60", "6", "Magnetometer [3]", "Magnetometer X, Y, Z"]
    ]

    # Add table
    rows = len(table_data)
    cols = len(table_data[0])
    left = Inches(5.2)
    top = Inches(1.9)
    width = Inches(4.3)
    height = Inches(3.2)

    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # Set column widths
    table.columns[0].width = Inches(0.6)
    table.columns[1].width = Inches(0.5)
    table.columns[2].width = Inches(1.4)
    table.columns[3].width = Inches(1.8)

    # Fill table
    for i, row_data in enumerate(table_data):
        for j, cell_text in enumerate(row_data):
            cell = table.cell(i, j)
            cell.text = cell_text

            # Format text
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(9)
                paragraph.alignment = PP_ALIGN.CENTER

                if i == 0:  # Header row
                    paragraph.font.bold = True
                    paragraph.font.size = Pt(10)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(220, 220, 220)

                # Highlight ADC row
                if i == 4:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(220, 255, 220)

            cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Note at bottom
    note_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(9), Inches(0.3))
    note_frame = note_box.text_frame
    p = note_frame.paragraphs[0]
    p.text = "The 8 ADC values are consecutive time samples from a single channel, representing 8ms of continuous data per packet."
    p.font.size = Pt(10)
    p.alignment = PP_ALIGN.CENTER

    # Save presentation
    prs.save('/media/eddlai/DATA/ExoPulse_firmware/uMyo_Blueprint_to_Beacon.pptx')
    print("PowerPoint presentation created successfully: uMyo_Blueprint_to_Beacon.pptx")

if __name__ == "__main__":
    create_presentation()
